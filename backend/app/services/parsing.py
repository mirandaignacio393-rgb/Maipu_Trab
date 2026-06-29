"""Extracción de texto desde distintos formatos de documento."""
from __future__ import annotations

import base64
import io
import re

import httpx

from ..clients import get_anthropic
from ..config import get_settings

_settings = get_settings()

_IMAGE_MEDIA = {
    "png": "image/png",
    "jpg": "image/jpeg",
    "jpeg": "image/jpeg",
    "gif": "image/gif",
    "webp": "image/webp",
}


def _ext(filename: str) -> str:
    return filename.rsplit(".", 1)[-1].lower() if "." in filename else ""


# ---------------------------------------------------------------------------
#  Extractores por formato
# ---------------------------------------------------------------------------
def _from_pdf(data: bytes) -> str:
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(data))
    return "\n\n".join((page.extract_text() or "") for page in reader.pages)


def _from_docx(data: bytes) -> str:
    from docx import Document

    doc = Document(io.BytesIO(data))
    parts = [p.text for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells]
            if any(cells):
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def _from_xlsx(data: bytes) -> str:
    from openpyxl import load_workbook

    wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    parts: list[str] = []
    for ws in wb.worksheets:
        parts.append(f"## Hoja: {ws.title}")
        for row in ws.iter_rows(values_only=True):
            values = [str(v) for v in row if v is not None]
            if values:
                parts.append(" | ".join(values))
    return "\n".join(parts)


def _from_pptx(data: bytes) -> str:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))
    parts: list[str] = []
    for idx, slide in enumerate(prs.slides, start=1):
        parts.append(f"## Diapositiva {idx}")
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                parts.append(shape.text_frame.text)
    return "\n".join(parts)


def _from_image(data: bytes, filename: str) -> str:
    """Usa Claude (visión) para describir/transcribir el contenido de la imagen."""
    media_type = _IMAGE_MEDIA.get(_ext(filename), "image/png")
    b64 = base64.standard_b64encode(data).decode("utf-8")
    client = get_anthropic()
    resp = client.messages.create(
        model=_settings.anthropic_model,
        max_tokens=2000,
        messages=[
            {
                "role": "user",
                "content": [
                    {
                        "type": "image",
                        "source": {"type": "base64", "media_type": media_type, "data": b64},
                    },
                    {
                        "type": "text",
                        "text": (
                            "Transcribí TODO el texto visible y describí el contenido relevante "
                            "de esta imagen (diagramas, tablas, pasos). Sé exhaustivo; esto se usará "
                            "como base de conocimiento."
                        ),
                    },
                ],
            }
        ],
    )
    return "".join(b.text for b in resp.content if b.type == "text")


def _from_link(url: str) -> str:
    """Descarga una página web y limpia el HTML de forma básica."""
    with httpx.Client(timeout=30, follow_redirects=True) as client:
        html = client.get(url).text
    html = re.sub(r"<(script|style)[^>]*>.*?</\1>", " ", html, flags=re.S | re.I)
    text = re.sub(r"<[^>]+>", " ", html)
    text = re.sub(r"\s+", " ", text)
    return text.strip()


# ---------------------------------------------------------------------------
#  Punto de entrada
# ---------------------------------------------------------------------------
def extract_text(
    *,
    doc_type: str,
    filename: str = "",
    data: bytes | None = None,
    source_url: str | None = None,
) -> str:
    """Devuelve el texto plano de un documento según su tipo."""
    dt = doc_type.lower()

    if dt in ("link",):
        if not source_url:
            raise ValueError("Falta source_url para un documento de tipo link")
        return _from_link(source_url)

    if dt in ("youtube",):
        # Transcripción automática de YouTube: pendiente (Fase posterior).
        # Por ahora guardamos la URL como referencia consultable.
        return f"Video de YouTube: {source_url}"

    if data is None:
        raise ValueError("Faltan los bytes del archivo")

    if dt == "pdf":
        return _from_pdf(data)
    if dt == "docx":
        return _from_docx(data)
    if dt == "xlsx":
        return _from_xlsx(data)
    if dt == "pptx":
        return _from_pptx(data)
    if dt == "image":
        return _from_image(data, filename)
    if dt in ("text", "txt", "md"):
        return data.decode("utf-8", errors="replace")

    raise ValueError(f"Tipo de documento no soportado: {doc_type}")
